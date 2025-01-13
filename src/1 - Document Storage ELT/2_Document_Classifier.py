import json
import logging
import os
import shutil  # Import shutil for moving files
import fitz  # Import PyMuPDF for PDF processing
import tabula  # Import tabula-py for table extraction
from docx import Document as DocxDocument  # Import python-docx for DOCX processing
from PIL import Image  # Import Pillow for image processing
from pptx import Presentation  # Import python-pptx for PPTX processing
import pandas as pd  # Import pandas for CSV and Excel processing
import pypandoc  # Import pypandoc for RTF processing
from bs4 import BeautifulSoup  # Import BeautifulSoup for HTML processing
import xml.etree.ElementTree as ET  # Import ElementTree for XML processing
import zipfile  # Import zipfile for ZIP processing
from odf.opendocument import load  # Import odfpy for ODT processing
from odf.text import P  # Import odfpy text elements
from odf.draw import Image as ODFImage  # Import odfpy image elements
from odf.table import Table  # Import odfpy table elements

import sys
import os  # Ensure os is imported

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from index_management import load_index, save_index, index_exists, cleanse_index
from file_operations import move_file, delete_file
from utils import generate_document_id, calculate_hash
from logging_setup import setup_logging, get_logger

# Import configuration
from config import ROOT_FOLDER, INDEX_FILE, LOGS_FOLDER, LOG_FILE, CLASSIFIED_FOLDER

# Define COMBINATIONS_FILE early
COMBINATIONS_FILE = os.path.join(ROOT_FOLDER, "data", "documents", "classification_combinations.json")

# Initialize logging using config
setup_logging(LOG_FILE)
logger = get_logger(__name__)

# Add a console handler to print only errors
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.ERROR)
console_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
logger.addHandler(console_handler)

# Ensure the directory for the log file exists
os.makedirs(LOGS_FOLDER, exist_ok=True)

# --- Configuration ---
ROOT_FOLDER = ROOT_FOLDER  # Already imported from config.py
INDEX_FILE = INDEX_FILE
CLASSIFIED_FOLDER = CLASSIFIED_FOLDER

# --- Logging Setup ---
os.makedirs(LOGS_FOLDER, exist_ok=True)

# Load the decision tree
def load_decision_tree(filepath):
    """Loads the decision tree from a JSON file."""
    decision_tree_path = os.path.join(os.path.dirname(__file__), filepath)
    try:
        with open(decision_tree_path, "r") as f:
            decision_tree = json.load(f)
        logging.info(f"Decision tree loaded from {decision_tree_path}")
        logging.debug(f"Decision tree content: {decision_tree}")  # Add debug logging
        return decision_tree
    except Exception as e:
        logging.error(f"Error loading decision tree: {e}")
        return None

# Classification functions
def classify_pdf_for_text(filepath):
    """Classifies a PDF document to check if it contains text."""
    try:
        doc = fitz.open(filepath)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text.strip():  # Check if the extracted text is non-empty
                return True
        return False
    except Exception as e:
        logging.error(f"Error classifying PDF for text: {e}")
        return False

def classify_pdf_for_images(filepath):
    """Classifies a PDF document to check if it contains images."""
    try:
        doc = fitz.open(filepath)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            images = page.get_images(full=True)
            if images:  # Check if any images are found
                return True
        return False
    except Exception as e:
        logging.error(f"Error classifying PDF for images: {e}")
        return False

def classify_pdf_for_tables(filepath):
    """Classifies a PDF document to check if it contains tables."""
    try:
        tables = tabula.read_pdf(filepath, pages='all', multiple_tables=True)
        if tables:  # Check if any tables are found
            return True
        return False
    except Exception as e:
        logging.error(f"Error classifying PDF for tables: {e}")
        return False

def classify_docx_document(filepath):
    """Classifies a DOCX document to analyze its structure."""
    try:
        doc = DocxDocument(filepath)
        has_text = False
        has_images = False
        has_tables = False

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                has_text = True
                break

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                has_images = True
                break

        if doc.tables:
            has_tables = True

        if has_text and has_images and has_tables:
            return "DOCX-Text-with-Images-and-Tables"
        elif has_text and has_images:
            return "DOCX-Text-with-Images"
        elif has_text and has_tables:
            return "DOCX-Text-with-Tables"
        elif has_text:
            return "DOCX-Text-Only"
        else:
            return "DOCX-Unknown"
    except Exception as e:
        logging.error(f"Error classifying DOCX document: {e}")
        return "DOCX-Unknown"

def classify_doc_document(filepath):
    """Classifies a DOC document to analyze its structure."""
    return classify_docx_document(filepath)

def classify_image_document(filepath):
    """Classifies an image document to analyze its content."""
    try:
        with Image.open(filepath) as img:
            format = img.format
            size = img.size
            mode = img.mode
            return f"Image-{format}-{size[0]}x{size[1]}-{mode}"
    except Exception as e:
        logging.error(f"Error classifying image document: {e}")
        return "Image-Unknown"

def classify_tiff_as_multipage(filepath):
    """Classifies a TIFF document to check if it is multipage."""
    try:
        with Image.open(filepath) as img:
            if getattr(img, "n_frames", 1) > 1:  # Check if the TIFF has multiple frames
                return True
        return False
    except Exception as e:
        logging.error(f"Error classifying TIFF as multipage: {e}")
        return False

def classify_pptx_document(filepath):
    """Classifies a PPTX document to analyze its structure."""
    try:
        presentation = Presentation(filepath)
        has_text = False
        has_images = False

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    has_text = True
                if shape.shape_type == 13:  # Shape type 13 corresponds to pictures
                    has_images = True

        if has_text and has_images:
            return "PPTX-Text-with-Images"
        elif has_text:
            return "PPTX-Text-Only"
        elif has_images:
            return "PPTX-Image-Only"
        else:
            return "PPTX-Unknown"
    except Exception as e:
        logging.error(f"Error classifying PPTX document: {e}")
        return "PPTX-Unknown"

def classify_text_document(filepath):
    """Classifies a TXT document to analyze its content."""
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            content = file.read()
            if content.strip():  # Check if the text content is non-empty
                return "Text-Only"
            else:
                return "Text-Empty"
    except Exception as e:
        logging.error(f"Error classifying TXT document: {e}")
        return "Text-Unknown"

def classify_csv_document(filepath):
    """Classifies a CSV document to analyze its structure."""
    try:
        df = pd.read_csv(filepath)
        if not df.empty:  # Check if the DataFrame is not empty
            return "Table-Only"
        else:
            return "Table-Empty"
    except Exception as e:
        logging.error(f"Error classifying CSV document: {e}")
        return "Table-Unknown"

def classify_excel_document(filepath):
    """Classifies an Excel document (XLS/XLSX) to analyze its structure."""
    try:
        if filepath.endswith('.xlsx'):
            df = pd.read_excel(filepath, engine='openpyxl')
        elif filepath.endswith('.xls'):
            df = pd.read_excel(filepath, engine='xlrd')
        else:
            return "Table-Unknown"

        if not df.empty:  # Check if the DataFrame is not empty
            return "Table-Only"
        else:
            return "Table-Empty"
    except Exception as e:
        logging.error(f"Error classifying Excel document: {e}")
        return "Table-Unknown"

def classify_rtf_document(filepath):
    """Classifies an RTF document to analyze its content."""
    try:
        content = pypandoc.convert_file(filepath, 'plain')
        if content.strip():  # Check if the text content is non-empty
            return "Text-Only"
        else:
            return "Text-Empty"
    except Exception as e:
        logging.error(f"Error classifying RTF document: {e}")
        return "Text-Unknown"

def classify_html_document(filepath):
    """Classifies an HTML document to analyze its content."""
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            soup = BeautifulSoup(file, "html.parser")
            has_text = bool(soup.get_text(strip=True))
            has_images = bool(soup.find_all("img"))
            if has_text and has_images:
                return "HTML-Text-with-Images"
            elif has_text:
                return "HTML-Text-Only"
            elif has_images:
                return "HTML-Image-Only"
            else:
                return "HTML-Unknown"
    except Exception as e:
        logging.error(f"Error classifying HTML document: {e}")
        return "HTML-Unknown"

def classify_xml_document(filepath):
    """Classifies an XML document to analyze its content."""
    try:
        tree = ET.parse(filepath)
        root = tree.getroot()
        has_text = bool(root.text.strip() if root.text else False)
        has_elements = bool(root.findall(".//*"))
        if has_text and has_elements:
            return "XML-Text-with-Elements"
        elif has_text:
            return "XML-Text-Only"
        elif has_elements:
            return "XML-Elements-Only"
        else:
            return "XML-Unknown"
    except Exception as e:
        logging.error(f"Error classifying XML document: {e}")
        return "XML-Unknown"

def classify_vtt_document(filepath):
    """Classifies a VTT document to analyze its content."""
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            content = file.read()
            if content.strip():
                return "Text-Only"
            else:
                return "Text-Empty"
    except Exception as e:
        logging.error(f"Error classifying VTT document: {e}")
        return "Text-Unknown"

def classify_zip_contents(filepath):
    """Classifies a ZIP document to analyze its content."""
    try:
        with zipfile.ZipFile(filepath, 'r') as zip_ref:
            if zip_ref.namelist():  # Check if the ZIP file contains any files
                return "Archive"
            else:
                return "Archive-Empty"
    except Exception as e:
        logging.error(f"Error classifying ZIP document: {e}")
        return "Archive-Unknown"

def classify_odt_document(filepath):
    """Classifies an ODT document to analyze its structure."""
    try:
        doc = load(filepath)
        has_text = False
        has_images = False
        has_tables = False

        for element in doc.getElementsByType(P):
            if element.text.strip():
                has_text = True
                break

        for element in doc.getElementsByType(ODFImage):
            has_images = True
            break

        for element in doc.getElementsByType(Table):
            has_tables = True
            break

        if has_text and has_images and has_tables:
            return "ODT-Text-with-Images-and-Tables"
        elif has_text and has_images:
            return "ODT-Text-with-Images"
        elif has_text and has_tables:
            return "ODT-Text-with-Tables"
        elif has_text:
            return "ODT-Text-Only"
        else:
            return "ODT-Unknown"
    except Exception as e:
        logging.error(f"Error classifying ODT document: {e}")
        return "ODT-Unknown"

# Modify the update_classification_combinations function to track counts
def update_classification_combinations(primary_classification, secondary_classifications):
    """Updates the classification combinations file with new combinations and counts."""
    try:
        if os.path.exists(COMBINATIONS_FILE):
            with open(COMBINATIONS_FILE, "r") as f:
                combinations = json.load(f)
        else:
            combinations = []

        # Sort secondary classifications for consistency
        sorted_secondary = sorted(secondary_classifications)

        # Check if the combination already exists and increment count
        for combo in combinations:
            if (combo["primary_classification"] == primary_classification and
                sorted(combo["secondary_classifications"]) == sorted_secondary):
                combo["count"] += 1
                logging.info(f"Incremented count for classification combination: {combo}")
                break
        else:
            new_combination = {
                "primary_classification": primary_classification,
                "secondary_classifications": sorted_secondary,
                "count": 1  # Initialize count
            }
            combinations.append(new_combination)
            logging.info(f"Added new classification combination: {new_combination}")

        with open(COMBINATIONS_FILE, "w") as f:
            json.dump(combinations, f, indent=4)
    except Exception as e:
        logging.error(f"Error updating classification combinations: {e}")

# Ensure the file extension is correctly formatted before looking it up in the decision tree.
def get_file_extension(filepath):
    """Returns the file extension for a given filepath."""
    return os.path.splitext(filepath)[1].lower().lstrip('.')

# Ensure that secondary classifications are standardized in classify_document
def classify_document(doc_id, doc_data, decision_tree, use_case_to_function, index):
    """Classifies a document based on the decision tree."""
    filepath = doc_data["filepath"]
    current_folder = os.path.basename(os.path.dirname(filepath))

    # Remove the forced 'return' when physically in Classified folder
    if "Classified" in os.path.dirname(filepath) and doc_data["status"] == "In_Processing":
        logging.info(
            f"Document {doc_id} is in Classified folder but index shows 'In_Processing'. "
            "Re-classifying and updating metadata instead of skipping."
        )
        # doc_data["status"] = "Classified"  # Remove or comment out this direct status change
        # ...continue classification process instead of returning...

    # Allow re-classification if file is in Duplicates folder but status is In_Processing
    if "Duplicates" in os.path.dirname(filepath) and doc_data["status"] == "In_Processing":
        logging.info(
            f"Document {doc_id} is physically in Duplicates folder but index shows 'In_Processing'. "
            "Overwriting duplicates and re-classifying."
        )
        # ...existing code to remove or overwrite duplicates if needed...
        # doc_data["status"] = "In_Processing"  # Ensure we keep it In_Processing for classification

    # ...existing code that re-classifies the document normally...
    # (Removed any early 'return' so classification isn't skipped)

    extension = get_file_extension(filepath)
    logging.debug(f"Classifying document ID: {doc_id}, Extension: {extension}")  # Add debug logging

    # Recalculate hash before processing
    doc_data["hash"] = calculate_hash(filepath)
    doc_id = doc_data["hash"]  # Ensure doc_id matches the file hash

    if extension not in decision_tree:
        outcome = f"No classification rule found for extension '{extension}'"
        logging.warning(f"{outcome} (document ID: {doc_id})")
        logging.debug(f"Available extensions in decision tree: {list(decision_tree.keys())}")  # Add debug logging
        return

    primary_classification = decision_tree[extension]["primary"]
    index[doc_id]["document_type"] = primary_classification
    index[doc_id]["metadata"]["document_type"] = primary_classification
    index[doc_id]["metadata"]["primary_classification"] = primary_classification

    # Remove old classification flags
    for flag in ["Text", "Images", "Tables"]:
        if flag in index[doc_id]["metadata"]:
            del index[doc_id]["metadata"][flag]

    secondary_use_cases = decision_tree[extension]["secondary"]["use_cases"]
    secondary_classifications = []

    index[doc_id]["metadata"]["secondary_classifications"] = []
    # Execute use cases in order
    for use_case in secondary_use_cases:
        if use_case in use_case_to_function:
            classification_function = use_case_to_function[use_case]
            result = classification_function(filepath)
            if result:
                # If it's True/False or a descriptive string, parse it into standard keys
                if isinstance(result, bool):
                    if result:
                        if use_case == "extract_text_from_pdf":
                            index[doc_id]["metadata"]["Text"] = "Yes"
                        elif use_case == "detect_images_in_pdf":
                            index[doc_id]["metadata"]["Images"] = "Yes"
                        elif use_case == "detect_tables_in_pdf":
                            index[doc_id]["metadata"]["Tables"] = "Yes"
                    else:
                        if use_case == "extract_text_from_pdf":
                            index[doc_id]["metadata"]["Text"] = "No"
                        elif use_case == "detect_images_in_pdf":
                            index[doc_id]["metadata"]["Images"] = "No"
                        elif use_case == "detect_tables_in_pdf":
                            index[doc_id]["metadata"]["Tables"] = "No"
                else:
                    # For strings like "Text-with-Images", parse into standard keys
                    parse_classification_result(result, index[doc_id]["metadata"])
                # Collect standardized secondary classifications
                for key in ["Text", "Images", "Tables"]:
                    if index[doc_id]["metadata"].get(key) == "Yes":
                        if key not in secondary_classifications:
                            secondary_classifications.append(key)
                logging.info(f"Use case '{use_case}' returned: {result}")
                index[doc_id]["metadata"]["secondary_classifications"].append(use_case)

    # Determine final document type
    document_type = determine_document_type(doc_id, index, primary_classification)
    index[doc_id]["document_type"] = document_type
    logging.info(f"Document {doc_id} classified as: {document_type}")

    # Log the classification outcome with detailed information
    logging.info(f"Completed classification for document ID: {doc_id}, Outcome: {document_type}")

    # Update classification combinations
    update_classification_combinations(primary_classification, secondary_classifications)
    logging.info(f"Classification combinations updated for document ID: {doc_id} with: {secondary_classifications}")

    # Move document to classified folder
    move_to_classified_folder(doc_id, doc_data, CLASSIFIED_FOLDER)
    logging.info(f"Document ID: {doc_id} moved to Classified folder as: {document_type}")

def parse_classification_result(result, metadata):
    """
    Parses any result string or boolean and sets metadata['Text'], metadata['Images'],
    or metadata['Tables'] = 'Yes' or 'No' based on the outcome.
    """
    # If result is a string describing content (e.g. "DOCX-Text-with-Images-and-Tables")
    if isinstance(result, str):
        # Convert to lowercase for easier checks
        lower_res = result.lower()
        if "text" in lower_res:
            metadata["Text"] = "Yes"
        if "images" in lower_res:
            metadata["Images"] = "Yes"
        if "tables" in lower_res:
            metadata["Tables"] = "Yes"
    # If result is a boolean from a detection function
    elif isinstance(result, bool):
        # Example usage: detect_images_in_pdf -> True => metadata["Images"] = "Yes"
        # This logic is set by calling code in classify_document.
        pass

def determine_document_type(doc_id, index, primary_classification):
    """Determines the final document type based on primary and secondary classifications."""
    metadata = index[doc_id]["metadata"]

    # Use standardized tags to build final doc type
    text_tag = "Text" if metadata.get("Text") == "Yes" else ""
    images_tag = "Images" if metadata.get("Images") == "Yes" else ""
    tables_tag = "Tables" if metadata.get("Tables") == "Yes" else ""
    parts = [part for part in [text_tag, images_tag, tables_tag] if part]

    if primary_classification == "PDF":
        if parts:
            return "PDF-" + "-".join(parts)
        else:
            return "PDF-Unknown"

    # Handling for DOCX documents
    elif primary_classification == "DOCX":
        if metadata.get("analyze_docx_structure") == "DOCX-Text-with-Images-and-Tables":
            return "Text-with-Images-and-Tables"
        elif metadata.get("analyze_docx_structure") == "DOCX-Text-with-Images":
            return "Text-with-Images"
        elif metadata.get("analyze_docx_structure") == "DOCX-Text-with-Tables":
            return "Text-with-Tables"
        elif metadata.get("analyze_docx_structure") == "DOCX-Text-Only":
            return "Text-Only"
        else:
            return "DOCX-Unknown"

    # Handling for DOC documents
    elif primary_classification == "DOC":
        return determine_document_type(doc_id, index, "DOCX")

    # Handling for image documents (JPEG, PNG, GIF, BMP)
    elif primary_classification in ("JPEG", "PNG", "GIF", "BMP"):
        return "Image-Only"

    # Handling for TIFF documents
    elif primary_classification == "TIFF":
        if metadata.get("classify_tiff_as_multipage"):
            return "Image-Only-Multi-Page"
        else:
            return "Image-Only-Single-Page"

    # Handling for PPTX documents
    elif primary_classification == "PPTX":
        if metadata.get("analyze_pptx_content") == "PPTX-Text-with-Images":
            return "Text-with-Images"
        elif metadata.get("analyze_pptx_content") == "PPTX-Text-Only":
            return "Text-Only"
        elif metadata.get("analyze_pptx_content") == "PPTX-Image-Only":
            return "Image-Only"
        else:
            return "PPTX-Unknown"

    # Handling for TXT, VTT, VRT, CSV, XLS, XLSX, RTF, HTML, XML, ZIP, ODT documents
    elif primary_classification in ("TXT", "VTT", "VRT", "CSV", "XLS", "XLSX", "RTF", "HTML", "XML", "ZIP", "ODT"):
        return "Text-Only"

    # Default to primary classification if no specific logic is defined
    else:
        return primary_classification

def move_to_classified_folder(doc_id, doc_data, classified_folder):
    """Moves the document to the classified folder based on its document type."""
    document_type = doc_data["document_type"]
    destination_folder = os.path.join(classified_folder, document_type)
    try:
        # Check if destination_folder exists before attempting to create
        if not os.path.exists(destination_folder):
            os.makedirs(destination_folder, exist_ok=True)
            logging.info(f"Created Classified subfolder: {destination_folder}")
        else:
            logging.debug(f"Classified subfolder already exists: {destination_folder}")

        new_filepath = os.path.join(destination_folder, os.path.basename(doc_data["filepath"]))
        if os.path.exists(new_filepath):
            logging.info(f"File {new_filepath} already exists, moving to Duplicates.")
            doc_data["status"] = "Duplicates"
            move_to_duplicates_folder(doc_id, doc_data)
            return

        shutil.move(doc_data["filepath"], new_filepath)
        doc_data["filepath"] = new_filepath
        doc_data["status"] = "Classified"
        logging.info(f"Moved document {doc_id} to {destination_folder} with type: {document_type}")
    except Exception as e:
        logging.error(f"Error moving document {doc_id} to {destination_folder}: {e}")

from datetime import datetime  # Import datetime for timestamping

def move_to_duplicates_folder(doc_id, doc_data):
    """Moves the duplicate document to the Duplicates folder.
    
    If the document is already in the Duplicates folder, appends a timestamp to the filename.
    """
    duplicates_folder = os.path.join(ROOT_FOLDER, "Duplicates")
    try:
        if not os.path.exists(duplicates_folder):
            os.makedirs(duplicates_folder, exist_ok=True)
            logging.info(f"Created Duplicates folder: {duplicates_folder}")
        else:
            logging.debug(f"Duplicates folder already exists: {duplicates_folder}")

        current_folder = os.path.dirname(doc_data["filepath"])
        filename = os.path.basename(doc_data["filepath"])

        if os.path.abspath(current_folder) == os.path.abspath(duplicates_folder):
            # Append timestamp to filename to indicate version
            name, ext = os.path.splitext(filename)
            timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
            new_filename = f"{name}_{timestamp}{ext}"
            logging.info(f"Appending timestamp to duplicate filename: {new_filename}")
        else:
            new_filename = filename

        new_filepath = os.path.join(duplicates_folder, new_filename)
        if os.path.exists(new_filepath):
            os.remove(new_filepath)  # Overwrite existing duplicates
        shutil.move(doc_data["filepath"], new_filepath)
        doc_data["filepath"] = new_filepath
        doc_data["status"] = "Duplicates"
        logging.info(f"Moved document {doc_id} to Duplicates folder as {new_filename}")
    except Exception as e:
        logging.error(f"Error moving document {doc_id} to Duplicates folder: {e}")

# --- Main Classification Logic ---

def classify_documents_sequentially(index, decision_tree, use_case_to_function):
    in_processing_folder = os.path.join(ROOT_FOLDER, "In_Processing")
    if not os.path.exists(in_processing_folder):
        logging.warning(f"In_Processing folder not found: {in_processing_folder}")
        return

    # For each file in In_Processing, hash it, then classify or re-classify
    for filename in os.listdir(in_processing_folder):
        filepath = os.path.join(in_processing_folder, filename)
        if not os.path.isfile(filepath):
            continue

        doc_hash = calculate_hash(filepath)
        if doc_hash not in index:
            index[doc_hash] = {
                "document_id": doc_hash,
                "filepath": filepath,
                "filename": filename,
                "extension": os.path.splitext(filename)[1],
                "status": "In_Processing",
                "metadata": {}
            }
        else:
            index[doc_hash]["filepath"] = filepath
            index[doc_hash]["status"] = "In_Processing"

        try:
            classify_document(doc_hash, index[doc_hash], decision_tree, use_case_to_function, index)
        except Exception as e:
            logging.error(f"Error classifying document {doc_hash}: {e}")
            index[doc_hash]["status"] = "Errors"
            save_index(INDEX_FILE, index)

    # Save index after processing all In_Processing folder documents
    save_index(INDEX_FILE, index)

def initialize_index(index_file):
    if not index_exists(index_file):
        index = {}
        save_index(index_file, index)
    else:
        index = load_index(index_file)
    return index

index = initialize_index(INDEX_FILE)

# --- Main Execution ---
if __name__ == "__main__":
    logging.info("Starting document classification process.")
    # Load the decision tree
    decision_tree = load_decision_tree("classification_decision_tree.json")
    if not decision_tree:
        logging.error("Decision tree loading failed. Exiting.")
        exit(1)  # Exit if decision tree loading failed

    # Mapping of use cases to classification functions
    use_case_to_function = {
        "extract_text_from_pdf": classify_pdf_for_text,
        "detect_images_in_pdf": classify_pdf_for_images,
        "detect_tables_in_pdf": classify_pdf_for_tables,
        "analyze_vtt_content": classify_vtt_document,  # Corrected mapping for VTT files
        "analyze_doc_structure": classify_doc_document,
        "analyze_image_content": classify_image_document,
        "analyze_tiff_for_multi_page": classify_tiff_as_multipage,
        "analyze_pptx_content": classify_pptx_document,
        "analyze_txt_content": classify_text_document,
        "analyze_csv_structure": classify_csv_document,
        "analyze_xlsx_content": classify_excel_document,
        "analyze_rtf_content": classify_rtf_document,
        "analyze_html_structure": classify_html_document,
        "analyze_xml_structure": classify_xml_document,
        "extract_and_classify_zip": classify_zip_contents,
        "analyze_odt_structure": classify_odt_document,
    }

    # Load the document index
    try:
        with open(INDEX_FILE, "r") as f:
            index = json.load(f)
        logging.info(f"Document index loaded from {INDEX_FILE}")
    except Exception as e:
        logging.error(f"Error loading document index: {e}")
        exit(1)  # Exit if document index loading failed

    # Ensure Classified folder exists
    if not os.path.exists(CLASSIFIED_FOLDER):
        try:
            os.makedirs(CLASSIFIED_FOLDER, exist_ok=True)
            logging.info(f"Created Classified folder: {CLASSIFIED_FOLDER}")
        except Exception as e:
            logging.error(f"Error creating Classified folder: {e}")
            exit(1)  # Exit if Classified folder creation failed
    else:
        logging.debug(f"Classified folder already exists: {CLASSIFIED_FOLDER}")

    # Classify documents in the In_Processing folder sequentially
    classify_documents_sequentially(index, decision_tree, use_case_to_function)

    # Save the updated document index
    try:
        with open(INDEX_FILE, "w") as f:
            json.dump(index, f, indent=4)
        logging.info("Updated document index saved successfully.")
    except Exception as e:
        logging.error(f"Error saving updated document index: {e}")

    logging.info("Document classification process completed.")