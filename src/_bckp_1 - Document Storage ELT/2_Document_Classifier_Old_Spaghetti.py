import json
import logging
import os
import shutil
import fitz
import tabula
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation
import pandas as pd
import pypandoc
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
import zipfile
from odf.opendocument import load
from odf.text import P
from odf.draw import Image as ODFImage
from odf.table import Table
import sys
import os

# need to break in parts by document types

# Ensure the project root is correctly identified
project_root = os.path.dirname(os.path.abspath(__file__))

# Construct paths relative to the project root
ROOT_FOLDER = project_root
INDEX_FILE = os.path.join(ROOT_FOLDER, "data", "documents", "index.json")
LOGS_FOLDER = os.path.join(ROOT_FOLDER, "data", "logs")
LOG_FILE = os.path.join(LOGS_FOLDER, "document_classifier.log")
CLASSIFIED_FOLDER = os.path.join(ROOT_FOLDER, "data", "documents", "Classified")
COMBINATIONS_FILE = os.path.join(ROOT_FOLDER, "data", "documents", "classification_combinations.json")
DUPLICATES_FOLDER = os.path.join(ROOT_FOLDER, "data", "documents", "Duplicates")

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from index_management import load_index, save_index, index_exists, cleanse_index
from file_operations import move_file, delete_file
from utils import generate_document_id, calculate_hash
from logging_setup import setup_logging, get_logger

# Initialize logging using config
setup_logging(LOG_FILE)
logger = get_logger(__name__)

# Add a console handler to print only errors
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.ERROR)
console_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
logger.addHandler(console_handler)

# Ensure the directory for the log file exists
os.makedirs(LOGS_FOLDER, exist_ok=True)

# --- Configuration ---
# Already imported and defined above

# --- Logging Setup ---
# Already defined above

# Load the decision tree
def load_decision_tree(filepath):
    """Loads the decision tree from a JSON file."""
    decision_tree_path = os.path.join(os.path.dirname(__file__), filepath)
    try:
        with open(decision_tree_path, "r") as f:
            decision_tree = json.load(f)
        logging.info(f"Decision tree loaded from {decision_tree_path}")
        logging.debug(f"Decision tree content: {decision_tree}")  # Add debug logging
        return decision_tree
    except Exception as e:
        logging.error(f"Error loading decision tree: {e}")
        return None

# Classification functions
def classify_pdf_for_text(filepath):
    """Classifies a PDF document to check if it contains text."""
    try:
        doc = fitz.open(filepath)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text.strip():  # Check if the extracted text is non-empty
                return {"Text": "Yes"}
        return {"Text": "No"}
    except Exception as e:
        logging.error(f"Error classifying PDF for text: {e}")
        return {"Text": "Unknown"}

def classify_pdf_for_images(filepath):
    """Classifies a PDF document to check if it contains images."""
    try:
        doc = fitz.open(filepath)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            images = page.get_images(full=True)
            if images:  # Check if any images are found
                return {"Images": "Yes"}
        return {"Images": "No"}
    except Exception as e:
        logging.error(f"Error classifying PDF for images: {e}")
        return {"Images": "Unknown"}

def classify_pdf_for_tables(filepath):
    """Classifies a PDF document to check if it contains tables."""
    try:
        tables = tabula.read_pdf(filepath, pages='all', multiple_tables=True)
        if tables:  # Check if any tables are found
            return {"Tables": "Yes"}
        return {"Tables": "No"}
    except Exception as e:
        logging.error(f"Error classifying PDF for tables: {e}")
        return {"Tables": "Unknown"}

def classify_docx_document(filepath):
    """Classifies a DOCX document to analyze its structure."""
    try:
        doc = DocxDocument(filepath)
        has_text = False
        has_images = False
        has_tables = False

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                has_text = True
                break

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                has_images = True
                break

        if doc.tables:
            has_tables = True

        result = {}
        if has_text:
            result["Text"] = "Yes"
        else:
            result["Text"] = "No"

        if has_images:
            result["Images"] = "Yes"
        else:
            result["Images"] = "No"

        if has_tables:
            result["Tables"] = "Yes"
        else:
            result["Tables"] = "No"

        return result
    except Exception as e:
        logging.error(f"Error classifying DOCX document: {e}")
        return {"Text": "Unknown", "Images": "Unknown", "Tables": "Unknown"}

def classify_doc_document(filepath):
    """Classifies a DOC document to analyze its structure."""
    return classify_docx_document(filepath)

def classify_image_document(filepath):
    """Classifies an image document to analyze its content."""
    try:
        with Image.open(filepath) as img:
            format = img.format
            size = img.size
            mode = img.mode
            return {"Format": format, "Size": f"{size[0]}x{size[1]}", "Mode": mode}
    except Exception as e:
        logging.error(f"Error classifying image document: {e}")
        return {"Format": "Unknown", "Size": "Unknown", "Mode": "Unknown"}

def classify_tiff_as_multipage(filepath):
    """Classifies a TIFF document to check if it is multipage."""
    try:
        with Image.open(filepath) as img:
            if getattr(img, "n_frames", 1) > 1:
                return {"Multipage": "Yes"}
            else:
                return {"Multipage": "No"}
    except Exception as e:
        logging.error(f"Error classifying TIFF as multipage: {e}")
        return {"Multipage": "Unknown"}

def classify_pptx_document(filepath):
    """Classifies a PPTX document to analyze its structure."""
    try:
        presentation = Presentation(filepath)
        has_text = False
        has_images = False

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    has_text = True
                if shape.shape_type == 13:
                    has_images = True

        result = {}
        if has_text:
            result["Text"] = "Yes"
        else:
            result["Text"] = "No"

        if has_images:
            result["Images"] = "Yes"
        else:
            result["Images"] = "No"

        return result
    except Exception as e:
        logging.error(f"Error classifying PPTX document: {e}")
        return {"Text": "Unknown", "Images": "Unknown"}

def classify_text_document(filepath):
    """Classifies a TXT document to analyze its content."""
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            content = file.read()
            if content.strip():
                return {"Text": "Yes"}
            else:
                return {"Text": "No"}
    except Exception as e:
        logging.error(f"Error classifying TXT document: {e}")
        return {"Text": "Unknown"}

def classify_csv_document(filepath):
    """Classifies a CSV document to analyze its structure."""
    try:
        df = pd.read_csv(filepath)
        if not df.empty:
            return {"Tables": "Yes"}
        else:
            return {"Tables": "No"}
    except Exception as e:
        logging.error(f"Error classifying CSV document: {e}")
        return {"Tables": "Unknown"}

def classify_excel_document(filepath):
    """Classifies an Excel document (XLS/XLSX) to analyze its structure."""
    try:
        if filepath.endswith('.xlsx'):
            df = pd.read_excel(filepath, engine='openpyxl')
        elif filepath.endswith('.xls'):
            df = pd.read_excel(filepath, engine='xlrd')
        else:
            return {"Tables": "Unknown"}

        if not df.empty:
            return {"Tables": "Yes"}
        else:
            return {"Tables": "No"}
    except Exception as e:
        logging.error(f"Error classifying Excel document: {e}")
        return {"Tables": "Unknown"}

def classify_rtf_document(filepath):
    """Classifies an RTF document to analyze its content."""
    try:
        content = pypandoc.convert_file(filepath, 'plain')
        if content.strip():
            return {"Text": "Yes"}
        else:
            return {"Text": "No"}
    except Exception as e:
        logging.error(f"Error classifying RTF document: {e}")
        return {"Text": "Unknown"}

def classify_html_document(filepath):
    """Classifies an HTML document to analyze its content."""
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            soup = BeautifulSoup(file, "html.parser")
            has_text = bool(soup.get_text(strip=True))
            has_images = bool(soup.find_all("img"))

            result = {}
            if has_text:
                result["Text"] = "Yes"
            else:
                result["Text"] = "No"

            if has_images:
                result["Images"] = "Yes"
            else:
                result["Images"] = "No"

            return result
    except Exception as e:
        logging.error(f"Error classifying HTML document: {e}")
        return {"Text": "Unknown", "Images": "Unknown"}

def classify_xml_document(filepath):
    """Classifies an XML document to analyze its content."""
    try:
        tree = ET.parse(filepath)
        root = tree.getroot()
        has_text = bool(root.text.strip() if root.text else False)
        has_elements = bool(root.findall(".//*"))

        result = {}
        if has_text:
            result["Text"] = "Yes"
        else:
            result["Text"] = "No"
        
        if has_elements:
            result["Elements"] = "Yes"
        else:
            result["Elements"] = "No"

        return result
    except Exception as e:
        logging.error(f"Error classifying XML document: {e}")
        return {"Text": "Unknown", "Elements": "Unknown"}

def classify_vtt_document(filepath):
    """Classifies a VTT document to analyze its content."""
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            content = file.read()
            if content.strip():
                return {"Text": "Yes"}
            else:
                return {"Text": "No"}
    except Exception as e:
        logging.error(f"Error classifying VTT document: {e}")
        return {"Text": "Unknown"}

def classify_zip_contents(filepath):
    """Classifies a ZIP document to analyze its content."""
    try:
        with zipfile.ZipFile(filepath, 'r') as zip_ref:
            if zip_ref.namelist():  # Check if the ZIP file contains any files
                return {"Type": "Archive"}
            else:
                return {"Type": "Archive-Empty"}
    except Exception as e:
        logging.error(f"Error classifying ZIP document: {e}")
        return {"Type": "Archive-Unknown"}

def classify_odt_document(filepath):
    """Classifies an ODT document to analyze its structure."""
    try:
        doc = load(filepath)
        has_text = False
        has_images = False
        has_tables = False

        for element in doc.getElementsByType(P):
            if element.text.strip():
                has_text = True
                break

        for element in doc.getElementsByType(ODFImage):
            has_images = True
            break

        for element in doc.getElementsByType(Table):
            has_tables = True
            break

        result = {}
        if has_text:
            result["Text"] = "Yes"
        else:
            result["Text"] = "No"

        if has_images:
            result["Images"] = "Yes"
        else:
            result["Images"] = "No"

        if has_tables:
            result["Tables"] = "Yes"
        else:
            result["Tables"] = "No"
        
        return result
    except Exception as e:
        logging.error(f"Error classifying ODT document: {e}")
        return {"Text": "Unknown", "Images": "Unknown", "Tables": "Unknown"}

def update_classification_combinations(primary_classification, secondary_classifications):
    """Updates the classification combinations file with new combinations and counts."""
    try:
        if os.path.exists(COMBINATIONS_FILE):
            with open(COMBINATIONS_FILE, "r") as f:
                combinations = json.load(f)
        else:
            combinations = []

        # Sort secondary classifications for consistency
        sorted_secondary = sorted(secondary_classifications)

        # Check if the combination already exists and increment count
        for combo in combinations:
            if (combo["primary_classification"] == primary_classification and
                sorted(combo["secondary_classifications"]) == sorted_secondary):
                combo["count"] += 1
                logging.info(f"Incremented count for classification combination: {combo}")
                break
        else:
            new_combination = {
                "primary_classification": primary_classification,
                "secondary_classifications": sorted_secondary,
                "count": 1  # Initialize count
            }
            combinations.append(new_combination)
            logging.info(f"Added new classification combination: {new_combination}")

        with open(COMBINATIONS_FILE, "w") as f:
            json.dump(combinations, f, indent=4)
    except Exception as e:
        logging.error(f"Error updating classification combinations: {e}")

def get_file_extension(filepath):
    """Returns the file extension for a given filepath."""
    return os.path.splitext(filepath)[1].lower().lstrip('.')

def classify_document(doc_id, doc_data, decision_tree, use_case_to_function, index):
    """Classifies a document based on the decision tree."""
    filepath = doc_data["filepath"]
    current_folder = os.path.basename(os.path.dirname(filepath))
    filename = os.path.basename(filepath)
    extension = get_file_extension(filepath)

    # Check if the file exists in the Classified folder
    document_type = doc_data.get("document_type", "Unknown")
    classified_filepath = os.path.join(CLASSIFIED_FOLDER, document_type, filename)

    if os.path.exists(classified_filepath):
        logging.info(f"Document {doc_id} with filename {filename} found in both In_Processing and Classified folders. Prioritizing re-processing.")
        # Reset metadata and status for re-processing
        doc_data["status"] = "In_Processing"
        doc_data["metadata"] = {}

    # Re-calculate hash before processing (if not already calculated)
    if "hash" not in doc_data:
        doc_data["hash"] = calculate_hash(filepath)

    # Use the existing extension or determine the primary classification based on the extension from the decision tree
    primary_classification = doc_data.get("document_type", "Unknown")
    if extension in decision_tree:
        primary_classification = decision_tree[extension]["primary"]

    # Update index with primary classification
    index[doc_id]["document_type"] = primary_classification

    index[doc_id]["metadata"]["primary_classification"] = primary_classification

    logging.debug(f"Classifying document ID: {doc_id}, Extension: {extension}, Primary Classification: {primary_classification}")

    secondary_classifications = []
    if extension in decision_tree:
        secondary_use_cases = decision_tree[extension]["secondary"]["use_cases"]
        for use_case in secondary_use_cases:
            if use_case in use_case_to_function:
                classification_function = use_case_to_function[use_case]
                result = classification_function(filepath)

                if isinstance(result, dict):
                    # Update metadata with results from classification function
                    doc_data["metadata"].update(result)

                    # Add detected features to secondary_classifications list
                    for key, value in result.items():
                        if value == "Yes":
                            secondary_classifications.append(key)
                else:
                    logging.error(f"Classification function {use_case} did not return a dictionary. Result: {result}")

                logging.info(f"Use case '{use_case}' returned: {result}")
            else:
                logging.warning(f"No function mapped for use case: {use_case}")

    # Determine the final document type
    document_type = determine_document_type(doc_id, index, primary_classification)
    doc_data["document_type"] = document_type
    index[doc_id]["document_type"] = document_type

    # Log the classification outcome
    logging.info(f"Document {doc_id} classified as: {document_type}")

    # Update classification combinations
    update_classification_combinations(primary_classification, secondary_classifications)

    # Update the index after successful classification
    try:
        if doc_data["status"] == "Classified":
            # Update the index with the new metadata and document type
            index[doc_id]["metadata"] = doc_data["metadata"]
            index[doc_id]["document_type"] = doc_data["document_type"]
            save_index(INDEX_FILE, index)
            logging.info(f"Index updated for document ID: {doc_id}")

            # Move to Duplicates folder
            move_to_duplicates_folder(doc_id, doc_data)
        else:
            logging.error(f"Classification failed for document {doc_id}.")
            if current_folder != "In_Processing":
                move_to_in_processing_folder(doc_id, doc_data)
    except Exception as e:
        logging.error(f"Error updating index or moving document {doc_id} to Duplicates: {e}")
        if current_folder != "In_Processing":
            move_to_in_processing_folder(doc_id, doc_data)

def move_to_duplicates_folder(doc_id, doc_data):
    """Moves the duplicate document to the Duplicates folder and overwrites if exists."""
    try:
        if not os.path.exists(DUPLICATES_FOLDER):
            os.makedirs(DUPLICATES_FOLDER, exist_ok=True)
            logging.info(f"Created Duplicates folder: {DUPLICATES_FOLDER}")
        else:
            logging.debug(f"Duplicates folder already exists: {DUPLICATES_FOLDER}")

        filename = os.path.basename(doc_data["filepath"])
        new_filepath = os.path.join(DUPLICATES_FOLDER, filename)

        # Overwrite if exists
        if os.path.exists(new_filepath):
            os.remove(new_filepath)
            logging.info(f"Overwriting existing file in Duplicates folder: {filename}")

        shutil.move(doc_data["filepath"], new_filepath)
        doc_data["filepath"] = new_filepath
        doc_data["status"] = "Duplicates"
        logging.info(f"Moved document {doc_id} to Duplicates folder as {filename}")
    except Exception as e:
        logging.error(f"Error moving document {doc_id} to Duplicates folder: {e}")

def move_to_in_processing_folder(doc_id, doc_data):
    """Moves the document back to the In_Processing folder."""
    in_processing_folder = os.path.join(ROOT_FOLDER, "In_Processing")
    try:
        if not os.path.exists(in_processing_folder):
            os.makedirs(in_processing_folder, exist_ok=True)
            logging.info(f"Created In_Processing folder: {in_processing_folder}")
        else:
            logging.debug(f"In_Processing folder already exists: {in_processing_folder}")

        new_filepath = os.path.join(in_processing_folder, os.path.basename(doc_data["filepath"]))
        shutil.move(doc_data["filepath"], new_filepath)
        doc_data["filepath"] = new_filepath
        logging.info(f"Moved document {doc_id} back to In_Processing folder.")
    except Exception as e:
        logging.error(f"Error moving document {doc_id} back to In_Processing folder: {e}")

def determine_document_type(doc_id, index, primary_classification):
    """Determines the final document type based on primary and secondary classifications."""
    metadata = index[doc_id]["metadata"]

    # Start with the primary classification
    document_type = primary_classification

    # Append secondary classifications if present
    if metadata.get("Text") == "Yes":
        document_type += "-Text"
    if metadata.get("Images") == "Yes":
        document_type += "-Images"
    if metadata.get("Tables") == "Yes":
        document_type += "-Tables"
    if metadata.get("Elements") == "Yes":
        document_type += "-Elements"

    return document_type

def classify_documents_sequentially(index, decision_tree, use_case_to_function):
    in_processing_folder = os.path.join(ROOT_FOLDER, "In_Processing")
    if not os.path.exists(in_processing_folder):
        logging.warning(f"In_Processing folder not found: {in_processing_folder}")
        return

    for filename in os.listdir(in_processing_folder):
        filepath = os.path.join(in_processing_folder, filename)
        if not os.path.isfile(filepath):
            continue

        doc_hash = calculate_hash(filepath)

        # Get document type based on extension
        extension = get_file_extension(filepath)

        # Fetch primary classification from decision tree based on extension
        primary_classification = "Unknown"  # Default to "Unknown"
        if extension in decision_tree:
            primary_classification = decision_tree[extension]["primary"]

        if doc_hash not in index:
            index[doc_hash] = {
                "document_id": doc_hash,
                "filepath": filepath,
                "filename": filename,
                "extension": extension,
                "document_type": primary_classification,  # Use primary classification
                "status": "In_Processing",
                "metadata": {},
                "hash": doc_hash, # Add hash to index
            }
        else:
            # If the document is already in the index, update the filepath and status
            index[doc_hash]["filepath"] = filepath
            index[doc_hash]["status"] = "In_Processing"

        # Classify the document
        try:
            classify_document(doc_hash, index[doc_hash], decision_tree, use_case_to_function, index)
        except Exception as e:
            logging.error(f"Error classifying document {doc_hash}: {e}")
            index[doc_hash]["status"] = "Errors"
            save_index(INDEX_FILE, index)

    # Save the index after processing all files
    save_index(INDEX_FILE, index)

def initialize_index(index_file):
    if not index_exists(index_file):
        index = {}
        save_index(index_file, index)
    else:
        index = load_index(index_file)
    return index

index = initialize_index(INDEX_FILE)

# --- Main Execution ---
if __name__ == "__main__":
    logging.info("Starting document classification process.")
    # Load the decision tree
    decision_tree = load_decision_tree("classification_decision_tree.json")
    if not decision_tree:
        logging.error("Decision tree loading failed. Exiting.")
        exit(1)  # Exit if decision tree loading failed

    # Mapping of use cases to classification functions
    use_case_to_function = {
        "extract_text_from_pdf": classify_pdf_for_text,
        "detect_images_in_pdf": classify_pdf_for_images,
        "detect_tables_in_pdf": classify_pdf_for_tables,
        "analyze_vtt_content": classify_vtt_document,
        "analyze_doc_structure": classify_doc_document,
        "analyze_docx_structure": classify_docx_document,
        "analyze_image_content": classify_image_document,
        "analyze_tiff_for_multi_page": classify_tiff_as_multipage,
        "analyze_pptx_content": classify_pptx_document,
        "analyze_txt_content": classify_text_document,
        "analyze_csv_structure": classify_csv_document,
        "analyze_xlsx_content": classify_excel_document,
        "analyze_rtf_content": classify_rtf_document,
        "analyze_html_structure": classify_html_document,
        "analyze_xml_structure": classify_xml_document,
        "extract_and_classify_zip": classify_zip_contents,
        "analyze_odt_structure": classify_odt_document,
    }

    # Load the document index
    try:
        with open(INDEX_FILE, "r") as f:
            index = json.load(f)
        logging.info(f"Document index loaded from {INDEX_FILE}")
    except Exception as e:
        logging.error(f"Error loading document index: {e}")
        exit(1)  # Exit if document index loading failed

    # Ensure Classified folder exists
    if not os.path.exists(CLASSIFIED_FOLDER):
        try:
            os.makedirs(CLASSIFIED_FOLDER, exist_ok=True)
            logging.info(f"Created Classified folder: {CLASSIFIED_FOLDER}")
        except Exception as e:
            logging.error(f"Error creating Classified folder: {e}")
            exit(1)  # Exit if Classified folder creation failed
    else:
        logging.debug(f"Classified folder already exists: {CLASSIFIED_FOLDER}")

    # Classify documents in the In_Processing folder sequentially
    classify_documents_sequentially(index, decision_tree, use_case_to_function)

    logging.info("Document classification process completed.")